import os
import sys
import json
from pathlib import Path
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from collections import defaultdict
from dotenv import load_dotenv
from rag.project_paths import get_capture_folder, get_capture_json_path

load_dotenv(dotenv_path=".env.local")

SECTION_LABELS = {
    "win theme": "win_themes",
    "win themes": "win_themes",
    "hot button": "hot_buttons",
    "hot buttons": "hot_buttons",
    "discriminator": "discriminators",
    "discriminators": "discriminators"
}

def normalize_heading(text: str) -> str:
    lower = text.lower().strip()
    for label, tag in SECTION_LABELS.items():
        if label in lower:
            return tag
    return None

def parse_docx(path: Path) -> list:
    print(f"📄 Parsing DOCX: {path.name}")
    doc = Document(path)
    chunks = []
    current_section = None
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        section = normalize_heading(text)
        if section:
            current_section = section
            continue
        if current_section:
            chunks.append({"section": current_section, "text": text, "source_file": path.name})
    return chunks

def parse_pptx(path: Path) -> list:
    print(f"📊 Parsing PPTX: {path.name}")
    prs = Presentation(path)
    chunks = []
    current_section = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            section = normalize_heading(text)
            if section:
                current_section = section
                continue
            if current_section:
                chunks.append({"section": current_section, "text": text, "source_file": path.name})
    return chunks

def parse_pdf(path: Path) -> list:
    print(f"📄 Parsing PDF: {path.name}")
    doc = fitz.open(str(path))
    chunks = []
    current_section = None
    for page in doc:
        for line in page.get_text("text").splitlines():
            text = line.strip()
            if not text:
                continue
            section = normalize_heading(text)
            if section:
                current_section = section
                continue
            if current_section:
                chunks.append({"section": current_section, "text": text, "source_file": path.name})
    return chunks

def parse_all_capture_files(folder_path: str) -> list:
    folder = Path(folder_path)
    capture_files = [f for f in folder.glob("capture*.*") if f.suffix.lower() in [".docx", ".pptx", ".ppt", ".pdf"]]
    if not capture_files:
        print("⚠️ No capture*.{docx,pptx,ppt,pdf} files found.")
        return []

    all_chunks = []
    section_counts = defaultdict(int)

    for file in capture_files:
        try:
            if file.suffix.lower() == ".docx":
                chunks = parse_docx(file)
            elif file.suffix.lower() in [".pptx", ".ppt"]:
                chunks = parse_pptx(file)
            elif file.suffix.lower() == ".pdf":
                chunks = parse_pdf(file)
            else:
                print(f"⚠️ Skipping unsupported file: {file.name}")
                continue

            print(f"   ➕ {len(chunks)} chunks extracted from {file.name}")
            for chunk in chunks:
                if not isinstance(chunk, dict) or "section" not in chunk or "text" not in chunk:
                    print(f"⚠️ Invalid chunk skipped in {file.name}: {chunk}")
                    continue
                section_counts[chunk["section"]] += 1
                all_chunks.append(chunk)

        except Exception as e:
            print(f"❌ Failed to parse {file.name}: {e}")

    print("\n📊 Section-wise chunk summary:")
    for section, count in section_counts.items():
        print(f"   - {section}: {count} chunks")

    return all_chunks

def save_chunks_json(data: list, output_path: str):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
    print(f"\n✅ Saved parsed capture chunks to: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python -m rag.parse_capture_file <portfolio> <opportunity>")
        sys.exit(1)

    portfolio = sys.argv[1]
    opportunity = sys.argv[2]
    folder = get_capture_folder(portfolio, opportunity)
    output_file = get_capture_json_path(portfolio, opportunity)

    chunks = parse_all_capture_files(folder)
    if chunks:
        save_chunks_json(chunks, output_file)
