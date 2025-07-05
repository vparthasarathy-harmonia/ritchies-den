# rag/extract_capture_themes.py

import os
import sys
import json
from pathlib import Path
from pptx import Presentation
from dotenv import load_dotenv
from langchain_community.document_loaders import UnstructuredFileLoader
from rag.project_paths import get_capture_folder, get_capture_json_path
from rag.llm_client_claude import invoke_claude

load_dotenv(dotenv_path=".env.local")

def extract_text_from_pptx(path: Path) -> str:
    try:
        print(f"📥 Parsing PPTX: {path.name}")
        prs = Presentation(path)
        slide_texts = []
        for i, slide in enumerate(prs.slides):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_lines.append(shape.text.strip())
            slide_text = "\n".join(slide_lines).strip()
            if slide_text:
                slide_texts.append(f"Slide {i+1}:\n{slide_text}")
        return "\n\n".join(slide_texts)
    except Exception as e:
        print(f"❌ Failed to parse {path.name}: {e}")
        return ""

def extract_text_from_other(path: Path) -> str:
    try:
        print(f"📄 Parsing with UnstructuredLoader: {path.name}")
        loader = UnstructuredFileLoader(str(path))
        docs = loader.load()
        return "\n\n".join([d.page_content for d in docs])
    except Exception as e:
        print(f"❌ Failed to load {path.name}: {e}")
        return ""

def parse_all_capture_files(folder_path: str) -> dict:
    folder = Path(folder_path)
    capture_files = [f for f in folder.glob("capture*.*") if f.suffix.lower() in [".pptx", ".ppt", ".pdf", ".docx", ".doc"]]

    if not capture_files:
        print("⚠️ No capture* files found.")
        return {}

    full_text = ""
    for file in capture_files:
        if file.suffix.lower() == ".pptx":
            full_text += extract_text_from_pptx(file) + "\n\n"
        else:
            full_text += extract_text_from_other(file) + "\n\n"

    if not full_text.strip():
        print("⚠️ No text extracted from capture files.")
        return {}

    return extract_win_themes_from_chunks(full_text.strip())

def extract_win_themes_from_chunks(full_text: str, chunk_size=3000, overlap=200) -> dict:
    chunks = []
    start = 0
    while start < len(full_text):
        end = min(start + chunk_size, len(full_text))
        chunks.append(full_text[start:end])
        start = end - overlap

    merged = {"pain_points": [], "win_themes": [], "differentiators": []}

    for i, chunk in enumerate(chunks):
        prompt = f"""
You are a proposal strategist.

Below is a portion of text from capture planning documents.
Your task is to extract and organize the following into structured bullet lists:

- Pain Points
- Win Themes
- Differentiators

Return your answer in this JSON format:
{{
  "pain_points": [...],
  "win_themes": [...],
  "differentiators": [...]
}}

---

Text:
{chunk}
"""
        try:
            response = invoke_claude(prompt)
            print(f"\n🧠 Claude response (chunk {i+1}): {response[:200]}...")
            json_start = response.find("{")
            parsed = json.loads(response[json_start:].strip())
            for key in merged:
                merged[key].extend(parsed.get(key, []))
        except Exception as e:
            print(f"❌ Error processing chunk {i+1}: {e}")

    for key in merged:
        merged[key] = sorted(set(merged[key]))
    return merged

def save_capture_json(data: dict, output_path: str):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
    print(f"✅ Saved parsed capture to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_capture_themes.py <portfolio> <opportunity>")
        sys.exit(1)

    portfolio = sys.argv[1]
    opportunity = sys.argv[2]

    folder = get_capture_folder(portfolio, opportunity)
    output_file = get_capture_json_path(portfolio, opportunity)

    structured = parse_all_capture_files(folder)
    if structured:
        save_capture_json(structured, output_file)
